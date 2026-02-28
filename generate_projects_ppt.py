from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    def add_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        tf = slide.shapes.placeholders[1].text_frame
        tf.text = bullet_points[0]
        for point in bullet_points[1:]:
            p = tf.add_paragraph()
            p.text = point

    # --- PROJECT 1: PACKAGE RISK ---
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "AI-Driven Logistics: Package Damage Risk Predictor"
    slide.placeholders[1].text = "Presented by challa-obulesh\nPredictive Analytics for Shipping Safety"

    add_slide("Project 1: The Problem", [
        "High operational costs due to damaged goods during transit.",
        "Lack of real-time risk visibility for logistics managers.",
        "Inefficient packaging decisions leading to waste."
    ])
    
    add_slide("Solution: Smart Risk Engine", [
        "Machine Learning models (Random Forest) predicting damage probability.",
        "Analysis of features: Handling quality, weather, route type, and transfers.",
        "Predictive interpretability using feature importance analysis."
    ])

    add_slide("Tech Stack & Deliverables", [
        "Python, Streamlit, Scikit-learn, Plotly.",
        "Interactive Web Interface (Streamlit).",
        "Automated Risk Reports and Visual Dashboards."
    ])

    # --- PROJECT 2: FITAI PRO ---
    # Title Slide fitAI
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "FitAI Pro: Multi-Modal Fitness Intelligence"
    slide.placeholders[1].text = "Next-Generation Health & Workout Companion"

    add_slide("Project 2: The Vision", [
        "Bridging the gap between manual tracking and automated health insights.",
        "A holistic platform for Nutrition, Exercise, and Progress tracking.",
        "Multi-modal inputs: Images (Meals) and Text (Workouts)."
    ])

    add_slide("Core Functional Modules", [
        "Meal Vision: AI-powered food detection and calorie calculation.",
        "Workout NLP: Convert natural language (e.g., '10 pushups') into logs.",
        "Health Predictor: Dynamic TDEE, BMI, and Macro goal setting.",
        "Visual Dashboard: Interactive progress charts and weekly planning."
    ])

    add_slide("Summary & Future Scope", [
        "Both projects deployed on GitHub & live via Streamlit Cloud.",
        "Scalable architecture prepared for real-world integration.",
        "Future: Wearable device integration and IoT tracking for logistics."
    ])

    prs.save('AI_Projects_Presentation.pptx')
    print("Success! 'AI_Projects_Presentation.pptx' has been created.")

if __name__ == "__main__":
    try:
        create_presentation()
    except Exception as e:
        print(f"Error: {e}")
